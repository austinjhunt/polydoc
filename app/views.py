'''
from django.shortcuts import render
from pptx import Presentation


# Create your views here.
from django.http import HttpResponse

def index(request):
	prs = Presentation('app/test.pptx')
	for slide in prs.slides:
		for shape in slide.shapes:
			if shape.has_text_frame:
				return HttpResponse(shape.text)
	return HttpResponse("No titles")
	''
	slides = prs.slides
	slide = slides[0]
	res = slide.title
	#return HttpResponse(slide.name)
	if res != None:
		return HttpResponse(res)
	else:
		return HttpResponse("No response")
	''
''' 
from http.client import HTTPResponse
from .models import DocumentContainer, Document
from .utils import DriveAPI
from .forms import DocumentUploadForm, UserLoginForm, UserCreateForm, DocumentContainerForm
from django.shortcuts import redirect,render 
from django.http import FileResponse, Http404
from django.views.generic import View, FormView, CreateView, UpdateView, DeleteView
from django.contrib.auth import logout, authenticate, login
from django.contrib.auth.views import LoginView
from django.contrib.auth.mixins import LoginRequiredMixin
from django.conf import settings 

class HomeView(View):
	def get(self, request): 
		return render(
			request, 
			template_name='index.html', 
			context={}
		)

class LogoutView(LoginRequiredMixin, View):
	def get(self, request):  
		logout(request)
		return redirect('home')

class RegisterView(FormView):
	form_class = UserCreateForm
	template_name = 'forms/register.html'
	success_url = '/' 
	def get(self, request):
		if request.user.is_authenticated:
			return redirect(self.success_url)
		else: 
			return render(request=request,template_name=self.template_name, context={'form': self.form_class()})
	def form_valid(self, form):   
		# This method is called when valid form data has been POSTed.
        # It should return an HttpResponse. 
		username = self.request.POST.get('username', None)
		password = self.request.POST.get('password1', None) 
		user = form.save()
		user = authenticate(username=username,password=password)
		login(self.request, user)
		return redirect(self.success_url) 

class LoginView(FormView):
	form_class = UserLoginForm
	template_name = 'forms/login.html'
	success_url = '/'
	def dispatch(self, request, *args, **kwargs):
		if request.user.is_authenticated:
			return redirect('home')
		else:
			return super(LoginView, self).dispatch(request, *args, **kwargs)

	def form_valid(self,form):
		user = form.save()
		user = authenticate(
			username=self.request.POST['username'],
			password=self.request.POST['password']
			)
		login(self.request, user)
		return super(LoginView,self).form_valid(form)

class ProfileView(LoginRequiredMixin, FormView):
	form_class = DocumentUploadForm
	def get(self, request): 
		user_document_containers = DocumentContainer.objects.filter(user=request.user)
		for doc_container in user_document_containers:
			setattr(
				doc_container, 
				'doc_count',
				Document.objects.filter(containers__in=[doc_container.id]).count()
			)
		return render(
			request=request,
			template_name='profile.html', 
			context = {
				'form': self.form_class(),
				'user_document_containers': user_document_containers,
				'user_documents': Document.objects.filter(user=request.user)}
		)
		
	def form_valid(self, form):
		## Called after POSTED data has been validated
		files = form.files # MultiValueDict with key "files"
		files = files.getlist('file')
		print(files)
		return render(
			request=self.request,
			template_name='profile.html', 
			context = {'form': form}
		)

class DocumentContainerCreateView(LoginRequiredMixin, CreateView):
	model = DocumentContainer
	form_class = DocumentContainerForm
	success_url = 'profile'
	login_url = 'login'
	template_name = 'forms/create-document-container.html'

	def form_valid(self, form): 
		# Create document container associated with this user 
		DocumentContainer(
			name=form.cleaned_data['name'],
			user=self.request.user
		).save()
		return redirect('profile')

class DocumentContainerUpdateView(LoginRequiredMixin, UpdateView):
	model = DocumentContainer
	form_class = DocumentContainerForm
	success_url = 'profile'
	login_url = 'login'
	template_name = 'forms/update-document-container.html'
	def form_valid(self, form):
		# get object id 
		object_id = self.request.POST.get('objectid', None)
		if object_id:
			print(f'Object ID is {object_id}')
			doc = DocumentContainer.objects.get(id=object_id)
			doc.name = form.cleaned_data['name'] 
			doc.save()
		return redirect('profile')

class DocumentContainerDeleteView(LoginRequiredMixin, DeleteView):
	model = DocumentContainer
	success_url = 'profile'
	login_url = 'login'
	template_name = 'forms/delete-document-container.html'
	def get_object(self, queryset=None):
		""" Hook to ensure object is owned by request.user. """
		obj = super(DocumentContainerDeleteView, self).get_object()
		if not obj.user == self.request.user:
			raise Http404
		return obj
	def get(self, request, pk):
		try:
			doc_container = DocumentContainer.objects.get(id=pk)
			if not doc_container.user == request.user:
				raise Http404
			related_docs = Document.objects.filter(containers__in=[doc_container])
			return render(
				request,
				self.template_name,
				context={ 
					'related_docs': related_docs,
					'object': doc_container
				}
			)
		except Exception as e:
			print(e)
			return redirect('home')

	def post(self, request, pk): 
		doc_container = DocumentContainer.objects.get(id=pk)
		if not doc_container.user == request.user:
			raise Http404
		Document.objects.filter(containers__in=[doc_container]).delete()
		doc_container.delete()
		return redirect('profile')

class DocumentCreateView(LoginRequiredMixin, View): 
	login_url = 'login'
	success_url = 'profile' 
	def get(self, request): 
		return render(
			request, 
			template_name='forms/create-document.html',
			context={
				'user_document_containers': DocumentContainer.objects.filter(user=request.user)
			} 
		)
	def post(self, request): 
		data = dict(request.POST).items() 
		document_containers = [
			key.split('-')[1] for key,val in data if key.startswith('document_container_id-')
		] 
		files = request.FILES # MultiValueDict with 'file' key
		files = files.getlist('file') # list of <InMemoryUploadedFile>
		save_to_path = f'{settings.MEDIA_ROOT}/documents/{request.user.username}'
		print(save_to_path)
		print(files)
		for f in files: 
			print(f) 
			clean_filename = f.__dict__['_name'].replace(' ','-').strip().lower()
			
			print(clean_filename)
			new_doc = Document(
				notes='',
				title=clean_filename,
				location=f'{save_to_path}/{clean_filename}',
				user=request.user, 
			)
			new_doc.file.save(clean_filename, f.file)
			new_doc.save()
			for container_id in document_containers:
				new_doc.containers.add(DocumentContainer.objects.get(id=container_id))
			new_doc.save()
		return redirect('profile')
		

class DocumentUpdateView(LoginRequiredMixin, CreateView):
	def get(self, request): 
		pass

class DocumentDeleteView(LoginRequiredMixin, CreateView):
	def get(self, request): 
		pass

def display_document(request):
    # Use credentials for whatever
    obj = DriveAPI()

    print("Attempting file download")

    # TODO: parameterize these so they aren't hardcoded

    # Change this if you want to change the id of the document that's being downloaded
    f_id = "1jN3sdlIlTy2W1Ce9OAxpie0vDJZUfOGK1MUKtlIlvuk"

    # This is the file that will be used to store the pdf version of the document on google drive
    f_name = "test_download.pdf"

    obj.FileDownload(f_id, f_name)
    try:
        return FileResponse(open(f_name, 'rb'), content_type='application/pdf')
    except FileNotFoundError:
        raise Http404()

